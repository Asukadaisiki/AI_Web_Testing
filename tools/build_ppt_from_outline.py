"""JSON-driven PPT builder — reads outline.json + design_brief.json, outputs v4.pptx.

Addresses layout issues from v3:
  - Dynamic font sizing prevents text overflow
  - Grid-based image/text layout prevents overlap
  - Multi-weight font hierarchy (Bold/Regular/Light)
  - Balanced content density per slide type
"""
from __future__ import annotations

import json
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.oxml import parse_xml
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "decks" / "ai-web-testing-review"
PIC = ROOT / "pic"
OUT = DECK / "build" / "ai-web-testing-v4.pptx"

# ── Font ──
FONT = "Microsoft YaHei"
FONT_LIGHT = "Microsoft YaHei Light"

# ── Colors ──
INK = RGBColor(17, 24, 39)
MUTED = RGBColor(75, 85, 99)
NAVY = RGBColor(17, 18, 42)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
BG_COLOR = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE_COLOR = RGBColor(229, 231, 235)
ZEBRA_EVEN = RGBColor(245, 247, 250)

# ── Canvas (16:9) ──
SLIDE_W = 13.333
SLIDE_H = 7.5
MX = 0.55
HEADER_Y = 0.35
SUB_Y = 0.85
BODY_TOP = 1.50
FOOTER_Y = 6.50
FOOTER_H = 0.45
CONTENT_W = SLIDE_W - 2 * MX  # 12.233
INNER_PAD = 0.15

# ── Font hierarchy ──
FONT_SIZES = {
    "h1": 28,
    "h1_long": 24,
    "h1_xlong": 20,
    "h2": 15,
    "body": 13,
    "body_small": 11,
    "card_title": 15,
    "card_body": 11,
    "table_header": 11,
    "table_cell": 10,
    "caption": 10,
    "footer": 10,
    "takeaway": 14,
    "cover_title": 32,
    "cover_sub": 15,
    "cover_tagline": 16,
    "cover_tech": 12,
}

# ── Image registry ──
PICS = {
    "sessions": PIC / "Snipaste_2026-05-13_19-47-45.png",
    "planning": PIC / "Snipaste_2026-05-13_19-48-02.png",
    "rerun": PIC / "Snipaste_2026-05-13_19-48-11.png",
    "cases": PIC / "Snipaste_2026-05-13_19-48-27.png",
    "edit": PIC / "Snipaste_2026-05-13_19-48-35.png",
    "report": PIC / "Snipaste_2026-05-13_19-48-46.png",
    "evidence": PIC / "Snipaste_2026-05-13_19-49-22.png",
}

ASSET_MAP = {
    "assets/Snipaste_2026-05-13_19-48-02.png": "planning",
    "assets/Snipaste_2026-05-13_19-48-11.png": "rerun",
    "assets/Snipaste_2026-05-13_19-48-27.png": "cases",
    "assets/Snipaste_2026-05-13_19-48-35.png": "edit",
    "assets/Snipaste_2026-05-13_19-49-22.png": "evidence",
    "assets/Snipaste_2026-05-13_19-48-46.png": "report",
    "assets/Snipaste_2026-05-13_19-47-45.png": "sessions",
}


# ══════════════════════════════════════════════
# Font helpers
# ══════════════════════════════════════════════

def east_asian_font(run, font_name=FONT):
    """Force East Asian font tag on a run for proper Chinese rendering."""
    rpr = run._r.get_or_add_rPr()
    rpr.set("lang", "zh-CN")
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rpr.find(tag, namespaces={"a": ns})
        if el is None:
            el = parse_xml(
                f'<{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                f'typeface="{font_name}"/>'
            )
            rpr.append(el)
        else:
            el.set("typeface", font_name)


def set_font(run, size, color=INK, bold=False, font_name=FONT):
    """Apply font properties to a run with EA font fix."""
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    east_asian_font(run, font_name)


# ══════════════════════════════════════════════
# Layout engine
# ══════════════════════════════════════════════

def estimate_text_size(text, font_size, max_w_in, max_h_in):
    """Return (fitted_font_size, needed_height_in) that fits within max_w x max_h.

    CJK char ≈ 0.75 * font_size points wide.  Line height ≈ 1.5 * font_size points.
    """
    usable_w_pt = (max_w_in - 2 * INNER_PAD) * 72
    chars_per_line = max(1, int(usable_w_pt / (font_size * 0.75)))

    def _count_lines(t):
        lines = 0
        for para in t.split("\n"):
            if not para:
                lines += 1
            else:
                lines += max(1, -(-len(para) // chars_per_line))
        return lines

    lines_needed = _count_lines(text)
    line_h_pt = font_size * 1.5
    needed_h_in = (lines_needed * line_h_pt) / 72 + 2 * INNER_PAD

    if needed_h_in <= max_h_in:
        return font_size, needed_h_in

    # Shrink font until it fits
    for s in range(font_size - 2, 7, -2):
        cpl = max(1, int(usable_w_pt / (s * 0.75)))
        # Re-count with new char-per-line
        ln = 0
        for para in text.split("\n"):
            if not para:
                ln += 1
            else:
                ln += max(1, -(-len(para) // cpl))
        need_h = (ln * s * 1.5) / 72 + 2 * INNER_PAD
        if need_h <= max_h_in:
            return s, need_h
    return 8, max_h_in


def fit_title(title, max_w=CONTENT_W):
    """Pick title font size based on character count."""
    if len(title) > 60:
        return FONT_SIZES["h1_xlong"]
    elif len(title) > 40:
        return FONT_SIZES["h1_long"]
    return FONT_SIZES["h1"]


def add_textbox(slide, text, x, y, w, h, size=13, color=INK, bold=False,
                align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                fill=None, border_color=None, font_name=FONT):
    """Create a textbox with proper EA font settings and inner padding."""
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(INNER_PAD)
    tf.margin_right = Inches(INNER_PAD)
    tf.margin_top = Inches(INNER_PAD)
    tf.margin_bottom = Inches(INNER_PAD)
    tf.vertical_anchor = valign
    tf.word_wrap = True

    for idx, line_text in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = line_text
        set_font(run, size, color, bold, font_name)
    return shape


# ── Slide decoration helpers ──

def add_bg(slide):
    """Full-slide background rectangle."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W), Inches(SLIDE_H)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = BG_COLOR
    rect.line.fill.background()


def add_header(slide, title, subtitle=None, page=None):
    """Standard page header with optional subtitle and page number."""
    page_reserve = 0.85 if page else 0
    title_w = CONTENT_W - page_reserve
    title_sz = fit_title(title, title_w)
    add_textbox(slide, title, MX, HEADER_Y, title_w, 0.50, title_sz, INK, True)
    if subtitle:
        add_textbox(slide, subtitle, MX + 0.02, SUB_Y, CONTENT_W - 1.0, 0.35,
                    FONT_SIZES["h2"], MUTED, font_name=FONT_LIGHT)
    if page:
        add_textbox(slide, page, SLIDE_W - MX - 0.80, HEADER_Y + 0.05, 0.80, 0.25,
                    10, MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, title, body, x, y, w, h, accent=BLUE):
    """Rounded card with colored title + auto-fit body text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE_COLOR
    shape.line.width = Pt(0.8)

    # Title at top of card
    add_textbox(slide, title, x + 0.15, y + 0.10, w - 0.30, 0.35,
                FONT_SIZES["card_title"], accent, True)

    # Body text fitted to remaining height
    body_h = h - 0.55
    if body_h < 0.3:
        return
    fit_sz, _ = estimate_text_size(body, FONT_SIZES["card_body"], w - 0.30, body_h)
    add_textbox(slide, body, x + 0.15, y + 0.50, w - 0.30, body_h,
                fit_sz, MUTED, font_name=FONT_LIGHT)


def add_takeaway(slide, text, y=FOOTER_Y):
    """Bottom takeaway bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(MX), Inches(y), Inches(CONTENT_W), Inches(FOOTER_H)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE_COLOR
    shape.line.width = Pt(0.8)
    add_textbox(slide, text, MX + 0.15, y + 0.05, CONTENT_W - 0.30, FOOTER_H - 0.10,
                FONT_SIZES["takeaway"], NAVY, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def add_fit_image(slide, image_path, x, y, w, h, border=True):
    """Place image maintaining aspect ratio, centered in bounding box.

    Returns (actual_x, actual_y, actual_w, actual_h) for overlap detection.
    """
    with Image.open(image_path) as im:
        iw, ih = im.size
    ratio = iw / ih
    box_ratio = w / h
    if box_ratio > ratio:
        ph = h
        pw = h * ratio
    else:
        pw = w
        ph = w / ratio
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    pic_shape = slide.shapes.add_picture(
        str(image_path), Inches(px), Inches(py), Inches(pw), Inches(ph)
    )
    if border:
        pic_shape.line.color.rgb = LINE_COLOR
        pic_shape.line.width = Pt(1.0)
    return (px, py, pw, ph)


# ══════════════════════════════════════════════
# Table renderer
# ══════════════════════════════════════════════

def _cell_fill(cell, color):
    """Set cell background color."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def add_table(slide, headers, rows, x, y, w, h, col_weights=None, caption=None):
    """Self-fitting table with adaptive column widths."""
    n_cols = len(headers)
    n_rows = len(rows)
    table_shape = slide.shapes.add_table(
        n_rows + 1, n_cols, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tbl = table_shape.table

    if col_weights:
        total_w = sum(col_weights)
        for i, cw in enumerate(col_weights):
            tbl.columns[i].width = Inches(w * cw / total_w)

    # Header row
    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(header)
        set_font(run, FONT_SIZES["table_header"], WHITE, True)
        _cell_fill(cell, NAVY)
        cell.text_frame.margin_left = Inches(0.08)
        cell.text_frame.margin_right = Inches(0.08)
        cell.text_frame.margin_top = Inches(0.05)
        cell.text_frame.margin_bottom = Inches(0.05)
        cell.text_frame.word_wrap = True

    # Data rows with zebra striping
    for r, row in enumerate(rows, 1):
        bg = WHITE if r % 2 else ZEBRA_EVEN
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            set_font(run, FONT_SIZES["table_cell"], INK)
            _cell_fill(cell, bg)
            cell.text_frame.margin_left = Inches(0.08)
            cell.text_frame.margin_right = Inches(0.08)
            cell.text_frame.margin_top = Inches(0.04)
            cell.text_frame.margin_bottom = Inches(0.04)
            cell.text_frame.word_wrap = True

    if caption:
        add_textbox(slide, caption, x, y + h + 0.05, w, 0.25,
                    FONT_SIZES["caption"], MUTED, font_name=FONT_LIGHT)
    return table_shape


# ══════════════════════════════════════════════
# Slide renderers
# ══════════════════════════════════════════════

def _resolve_image(assets):
    """Resolve asset hero_image to a PICS key."""
    if not assets:
        return None
    hero = assets.get("hero_image", "")
    return ASSET_MAP.get(hero)


def _parse_highlight(text):
    """Parse 'Title: Body' or 'Title：Body' into (title, body)."""
    for sep in (": ", "："):
        if sep in text:
            parts = text.split(sep, 1)
            return parts[0], parts[1]
    return text[:20], text


def render_title(prs, slide_data, page_num=0, total=0):
    """Cover slide."""
    _ = page_num, total  # cover has no page number
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")
    assets = slide_data.get("assets", {})
    notes = slide_data.get("notes", "")

    add_textbox(slide, title, MX, 0.88, 6.0, 1.30, FONT_SIZES["cover_title"], INK, True)
    add_textbox(slide, subtitle, MX + 0.05, 2.40, 6.0, 0.55,
                FONT_SIZES["cover_sub"], MUTED, font_name=FONT_LIGHT)
    if notes:
        add_textbox(slide, notes, MX + 0.05, 3.20, 5.80, 0.80,
                    FONT_SIZES["cover_tagline"], NAVY, True, fill=WHITE, border_color=LINE_COLOR)

    # Hero image
    hero_key = _resolve_image(assets)
    if hero_key and hero_key in PICS:
        add_fit_image(slide, PICS[hero_key], 6.85, 0.82, 5.80, 3.20)

    # Second image (report) if planning is the hero
    if hero_key == "planning" and "report" in PICS:
        add_fit_image(slide, PICS["report"], 6.85, 4.20, 5.80, 2.40)

    add_textbox(slide, "FastAPI  ·  React/Vite  ·  Playwright  ·  AI DSL  ·  DOM Locator",
                MX, 6.55, 6.0, 0.35, FONT_SIZES["cover_tech"], MUTED, font_name=FONT_LIGHT)

    # GitHub repo URL
    repo_url = slide_data.get("repo_url", "")
    if repo_url:
        add_textbox(slide, f"GitHub: {repo_url}",
                    MX, 6.90, 6.0, 0.30, 11, BLUE, font_name=FONT_LIGHT)
    return slide


def render_split(prs, slide_data, page_num, total):
    """Split layout: highlights grid + summary callout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, slide_data.get("title", ""), slide_data.get("subtitle", ""),
               f"{page_num}/{total}")

    highlights = slide_data.get("highlights", [])
    body_text = slide_data.get("body", "")

    n = len(highlights)
    if n == 0 and body_text:
        fit_sz, _ = estimate_text_size(body_text, FONT_SIZES["body"], CONTENT_W, 4.5)
        add_textbox(slide, body_text, MX, BODY_TOP, CONTENT_W, 4.5, fit_sz, INK)
    elif n > 0:
        cols = 2
        rows_per = -(-n // cols)  # ceil
        card_w = (CONTENT_W - 0.30) / 2
        card_h = min(1.55, 3.8 / rows_per)
        colors = [BLUE, GREEN, AMBER, RED]
        for i, hl in enumerate(highlights):
            col = i // rows_per
            row = i % rows_per
            cx = MX + col * (card_w + 0.30)
            cy = BODY_TOP + row * (card_h + 0.20)
            ct, cb = _parse_highlight(hl)
            add_card(slide, ct, cb, cx, cy, card_w, card_h, colors[i % len(colors)])

    callout = slide_data.get("summary_callout", "")
    if callout:
        add_takeaway(slide, callout)
    return slide


def render_image_sidebar(prs, slide_data, page_num, total):
    """Image + sidebar cards layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, slide_data.get("title", ""), slide_data.get("subtitle", ""),
               f"{page_num}/{total}")

    image_side = slide_data.get("image_side", "left")
    sections = slide_data.get("sidebar_sections", [])
    assets = slide_data.get("assets", {})
    callout = slide_data.get("summary_callout", "")

    img_w = CONTENT_W * 0.53
    img_h = 4.45
    gap = 0.30
    text_w = CONTENT_W - img_w - gap
    img_x = MX if image_side == "left" else MX + CONTENT_W - img_w
    text_x = MX + img_w + gap if image_side == "left" else MX
    img_y = BODY_TOP

    hero_key = _resolve_image(assets)
    if hero_key and hero_key in PICS:
        add_fit_image(slide, PICS[hero_key], img_x, img_y, img_w, img_h)

    card_colors = [BLUE, GREEN, AMBER]
    n_sec = len(sections)
    if n_sec > 0:
        card_h = min(1.15, (img_h - (n_sec - 1) * 0.12) / n_sec)
        for i, sec in enumerate(sections):
            cy = img_y + i * (card_h + 0.12)
            add_card(slide, sec.get("title", ""), sec.get("body", ""),
                     text_x, cy, text_w, card_h, card_colors[i % 3])

    if callout:
        add_takeaway(slide, callout)
    return slide


def render_table_slide(prs, slide_data, page_num, total):
    """Table layout slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, slide_data.get("title", ""), slide_data.get("subtitle", ""),
               f"{page_num}/{total}")

    table_data = slide_data.get("table", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    col_weights = table_data.get("column_weights")
    caption = table_data.get("caption", "")

    table_h = min(len(rows) * 0.55 + 0.45, 5.0)
    add_table(slide, headers, rows, MX, BODY_TOP, CONTENT_W, table_h,
              col_weights, caption)

    callout = slide_data.get("summary_callout", "")
    if callout:
        add_takeaway(slide, callout)
    return slide


def render_comparison_2col(prs, slide_data, page_num, total):
    """Two-column comparison with verdict."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, slide_data.get("title", ""), slide_data.get("subtitle", ""),
               f"{page_num}/{total}")

    left = slide_data.get("left", {})
    right = slide_data.get("right", {})
    verdict = slide_data.get("verdict", "")

    col_w = (CONTENT_W - 0.30) / 2
    col_h = 3.80

    add_card(slide, left.get("title", ""),
             "\n".join(left.get("body", [])),
             MX, BODY_TOP, col_w, col_h, RED)

    add_card(slide, right.get("title", ""),
             "\n".join(right.get("body", [])),
             MX + col_w + 0.30, BODY_TOP, col_w, col_h, BLUE)

    if verdict:
        add_textbox(slide, verdict, MX, BODY_TOP + col_h + 0.20, CONTENT_W, 0.60,
                    FONT_SIZES["takeaway"], NAVY, True, PP_ALIGN.CENTER,
                    MSO_ANCHOR.MIDDLE, fill=WHITE, border_color=LINE_COLOR)

    callout = slide_data.get("summary_callout", "")
    if callout:
        add_takeaway(slide, callout)
    return slide


def render_timeline(prs, slide_data, page_num, total):
    """Timeline with milestone cards in horizontal bands."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, slide_data.get("title", ""), slide_data.get("subtitle", ""),
               f"{page_num}/{total}")

    milestones = slide_data.get("milestones", [])
    n = len(milestones)
    if n == 0:
        return slide

    card_w = (CONTENT_W - (n - 1) * 0.25) / n
    card_h = 3.80
    colors = [BLUE, GREEN, AMBER, RED]

    for i, ms in enumerate(milestones):
        cx = MX + i * (card_w + 0.25)
        cy = BODY_TOP

        add_textbox(slide, ms.get("label", ""),
                    cx, cy, card_w, 0.35, FONT_SIZES["card_title"], WHITE, True,
                    PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, fill=colors[i % len(colors)])

        add_card(slide, ms.get("title", ""), ms.get("body", ""),
                 cx, cy + 0.45, card_w, card_h - 0.45, colors[i % len(colors)])

    callout = slide_data.get("summary_callout", "")
    if callout:
        add_takeaway(slide, callout)
    return slide


def render_standard(prs, slide_data, page_num, total):
    """Fallback renderer: title + body."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, slide_data.get("title", ""), slide_data.get("subtitle", ""),
               f"{page_num}/{total}")

    body = slide_data.get("body", "")
    if body:
        fit_sz, _ = estimate_text_size(body, FONT_SIZES["body"], CONTENT_W, 4.5)
        add_textbox(slide, body, MX, BODY_TOP, CONTENT_W, 4.5, fit_sz, INK)

    callout = slide_data.get("summary_callout", "")
    if callout:
        add_takeaway(slide, callout)
    return slide


RENDERERS = {
    "title": render_title,
    "split": render_split,
    "image-sidebar": render_image_sidebar,
    "table": render_table_slide,
    "comparison-2col": render_comparison_2col,
    "timeline": render_timeline,
    "standard": render_standard,
}


# ══════════════════════════════════════════════
# Main build
# ══════════════════════════════════════════════

def build():
    with open(DECK / "outline.json", "r", encoding="utf-8") as f:
        outline = json.load(f)

    slides = outline.get("slides", [])
    total = len(slides)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    OUT.parent.mkdir(parents=True, exist_ok=True)

    for i, slide_data in enumerate(slides):
        variant = slide_data.get("variant", slide_data.get("type", "standard"))
        # "content" type in outline maps to variant field
        if variant == "content":
            variant = slide_data.get("variant", "standard")
        renderer = RENDERERS.get(variant, render_standard)
        renderer(prs, slide_data, i + 1, total)
        print(f"  Slide {i+1}/{total}: {slide_data.get('title', '')[:60]} [{variant}]")

    prs.save(str(OUT))
    print(f"\nSaved: {OUT}")


if __name__ == "__main__":
    build()
