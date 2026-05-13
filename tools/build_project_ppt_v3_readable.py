from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from PIL import Image


ROOT = Path(__file__).resolve().parents[1]
PIC = ROOT / "pic"
OUT = ROOT / "AI_Web_Testing_项目展示_评审阅读版_v3_中文优化.pptx"

FONT = "Microsoft YaHei"
INK = RGBColor(17, 24, 39)
MUTED = RGBColor(75, 85, 99)
NAVY = RGBColor(17, 18, 42)
BLUE = RGBColor(37, 99, 235)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
BG = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(229, 231, 235)


def pic(name: str) -> Path:
    return PIC / name


PICS = {
    "sessions": pic("Snipaste_2026-05-13_19-47-45.png"),
    "planning": pic("Snipaste_2026-05-13_19-48-02.png"),
    "rerun": pic("Snipaste_2026-05-13_19-48-11.png"),
    "cases": pic("Snipaste_2026-05-13_19-48-27.png"),
    "edit": pic("Snipaste_2026-05-13_19-48-35.png"),
    "report": pic("Snipaste_2026-05-13_19-48-46.png"),
    "evidence": pic("Snipaste_2026-05-13_19-49-22.png"),
}


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.strip("#")
    return RGBColor(int(hex_value[:2], 16), int(hex_value[2:4], 16), int(hex_value[4:], 16))


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor = LINE, width: float = 0.8) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 16,
             color: RGBColor = INK, bold: bool = False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, fill: RGBColor | None = None,
             line: RGBColor | None = None, margin: float = 0.08):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        set_fill(shape, fill)
    else:
        shape.fill.background()
    if line:
        set_line(shape, line)
    else:
        shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    lines = text.split("\n")
    for idx, line_text in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line_text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        # Explicit East Asian font for Chinese rendering.
        rpr = run._r.get_or_add_rPr()
        rpr.set("lang", "zh-CN")
        for tag in ("a:latin", "a:ea", "a:cs"):
            el = rpr.find(tag, namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
            if el is None:
                from pptx.oxml import parse_xml
                el = parse_xml(f'<{tag} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{FONT}"/>')
                rpr.append(el)
            else:
                el.set("typeface", FONT)
    return shape


def add_bg(slide):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    set_fill(rect, BG)
    rect.line.fill.background()


def add_header(slide, title: str, subtitle: str | None = None, page: str | None = None):
    add_text(slide, title, 0.55, 0.35, 10.6, 0.45, 23, INK, True)
    if subtitle:
        add_text(slide, subtitle, 0.57, 0.85, 10.9, 0.35, 11, MUTED)
    if page:
        add_text(slide, page, 12.0, 0.42, 0.7, 0.25, 10, MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, title: str, body: str, x: float, y: float, w: float, h: float,
             accent: RGBColor = BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, WHITE)
    set_line(shape, LINE)
    add_text(slide, title, x + 0.18, y + 0.14, w - 0.30, 0.32, 15, accent, True)
    add_text(slide, body, x + 0.22, y + 0.55, w - 0.35, h - 0.65, 11, MUTED)


def add_takeaway(slide, text: str, y: float = 6.58):
    add_text(slide, text, 0.72, y, 11.9, 0.45, 14, NAVY, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, WHITE, LINE)


def add_fit_image(slide, image_path: Path, x: float, y: float, w: float, h: float, border=True):
    with Image.open(image_path) as im:
        iw, ih = im.size
    ratio = iw / ih
    box_ratio = w / h
    if box_ratio > ratio:
        ph = h
        pw = h * ratio
    else:
        pw = w
        ph = w / ratio
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    pic_shape = slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), Inches(pw), Inches(ph))
    if border:
        set_line(pic_shape, LINE, 1.0)
    return pic_shape


def add_table(slide, headers, rows, x, y, w, h, widths=None, font_size=11):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = table_shape.table
    if widths:
        total = sum(widths)
        for i, val in enumerate(widths):
            tbl.columns[i].width = Inches(w * val / total)
    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = header
        set_fill(cell, NAVY)
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.text_frame.paragraphs[0].runs[0].font.name = FONT
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(font_size)
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            fill = RGBColor(255, 255, 255) if r % 2 else RGBColor(245, 247, 250)
            set_fill(cell, fill)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = INK
    return table_shape


def add_image_slide(prs, page: int, title: str, subtitle: str, image_key: str, bullets: list[tuple[str, str]],
                    takeaway: str, image_side: str = "left"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, title, subtitle, f"{page}/15")
    if image_side == "left":
        img_x, text_x = 0.62, 8.25
    else:
        img_x, text_x = 5.05, 0.72
    add_fit_image(slide, PICS[image_key], img_x, 1.55, 7.55, 4.35)
    colors = [BLUE, GREEN, AMBER]
    for idx, (bt, bb) in enumerate(bullets):
        add_card(slide, bt, bb, text_x, 1.55 + idx * 1.35, 4.25, 1.05, colors[idx % 3])
    add_takeaway(slide, takeaway)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text(slide, "AI 增强型 Web UI\n自动化测试平台", 0.75, 0.88, 5.9, 1.25, 32, INK, True)
    add_text(slide, "评审自主阅读版｜自然语言规划 · 结构化 DSL · 稳定执行 · 证据报告", 0.80, 2.45, 5.9, 0.55, 15, MUTED)
    add_text(slide, "项目定位：不是通用浏览器 Agent，而是可治理、可验证、可复用的测试工程平台。", 0.80, 3.25, 5.55, 0.80, 16, NAVY, True, fill=WHITE, line=LINE)
    add_fit_image(slide, PICS["planning"], 6.78, 0.82, 5.85, 3.25)
    add_fit_image(slide, PICS["report"], 6.78, 4.25, 5.85, 2.35)
    add_text(slide, "FastAPI  ·  React/Vite  ·  Playwright  ·  AI DSL  ·  DOM Locator", 0.82, 6.55, 5.75, 0.35, 12, MUTED)

    # 2 Guide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "阅读导览：这份材料回答四个评审问题", "每页都保留可独立阅读的结论、解释和项目截图证据。", "2/15")
    cards = [
        ("创意背景", "为什么需要 AI 增强测试平台，以及为什么不能让 AI 直接绕过测试工程治理。", BLUE),
        ("技术工艺", "自然语言先转 DSL，再由后端 Runner 执行；定位采用 DOM 优先、视觉兜底。", GREEN),
        ("应用前景", "面向 Web 冒烟测试、回归测试、测试资产沉淀和测试服务交付。", AMBER),
        ("成本回报", "AI 按需调用，常规执行走 Playwright；收益来自减少重复回归和排障时间。", RED),
    ]
    for i, (t, b, c) in enumerate(cards):
        add_card(slide, t, b, 0.82 + (i % 2) * 6.15, 1.65 + (i // 2) * 1.85, 5.55, 1.35, c)
    add_takeaway(slide, "建议阅读路径：先理解项目价值，再看功能截图与技术链路，最后判断应用场景和投入回报。")

    # 3 Background
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "创意背景：AI 自动化测试需要“可控”，不只是“能点页面”", "正式测试需要可审查、可复跑、可追踪，这与通用浏览器 Agent 的目标不同。", "3/15")
    add_table(slide, ["评审问题", "传统做法的不足", "本项目的应对"], [
        ["脚本编写成本", "测试人员需要手写 Playwright 步骤、选择器和断言，起步门槛高。", "AI 先根据业务目标生成 DSL 草案，再由人工审查和保存为正式用例。"],
        ["元素定位稳定性", "CSS/XPath 容易随页面改版、动态 DOM、同名按钮而失效。", "DOM 语义定位优先，结合上下文定位、AI visual 兜底和人工修正沉淀。"],
        ["AI 执行可信度", "纯 AI Agent 可能省略步骤、误改断言，且难以版本化和复跑。", "AI 不直接作为正式执行源，所有执行必须经过结构化 DSL 与后端 Runner。"],
        ["失败排查效率", "传统报告常停留在 pass/fail，缺少定位路径和页面证据。", "每一步沉淀截图、URL、DOM、定位候选、最终命中和失败原因。"]
    ], 0.75, 1.55, 11.85, 4.35, [0.18, 0.41, 0.41], 11)
    add_takeaway(slide, "核心判断：测试平台要把 AI 放进工程治理链路，而不是让 AI 替代治理链路。")

    # 4 Goal table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "项目目标：形成从需求到报告的自动化测试闭环", "目标不是生成一次性脚本，而是沉淀可复用的测试资产。", "4/15")
    headers = ["阶段", "平台处理", "输出"]
    rows = [
        ["需求收集", "AI Planning 提取业务目标、入口 URL、核心流程、关键断言和测试数据", "结构化测试目标"],
        ["用例生成", "LLM 根据页面探索结果生成 DSL 草案，并经过 schema 校验与归一化", "可审查测试用例"],
        ["正式执行", "后端 Playwright Runner 执行 DSL，前端只负责触发和展示", "可信执行结果"],
        ["证据沉淀", "记录截图、DOM、日志、定位候选、最终命中和失败原因", "步骤级报告"],
        ["持续优化", "根据失败记录、人工修正和回归结果改进定位与生成规则", "更稳定的测试资产"]
    ]
    add_table(slide, headers, rows, 0.75, 1.55, 11.85, 4.35, [0.18, 0.56, 0.26], 12)
    add_takeaway(slide, "边界原则：AI 生成必须经过 DSL 校验；正式结果只以后端 Runner 的结构化执行数据为准。")

    # Image slides
    add_image_slide(prs, 5, "功能证据 1：AI 规划工作台", "需求结构化、工具调用、测试点生成和 DSL 草案在同一页面闭环。", "planning", [
        ("可读信息", "评委可直接看到被测系统、业务目标、入口 URL、核心流程、关键断言和测试账号。"),
        ("工具调用", "Agent 调用 create_project、explore_flow 等工具采集页面信息，而不是凭空生成脚本。"),
        ("工程意义", "AI 规划被放在平台流程里，后续草案仍需保存为用例并由后端执行。")
    ], "该页证明项目已具备“AI 对话式规划 + 工具调用 + 测试草案生成”的可视化闭环。")

    add_image_slide(prs, 6, "功能证据 2：复测与失败迭代闭环", "自动化测试平台真正的价值在于持续复测、分析失败并沉淀改进。", "rerun", [
        ("失败不是终点", "同一测试从 failed 到 passed 的迭代记录，说明平台支持真实调试过程。"),
        ("问题可分类", "失败可区分 assertion、runner 等类型，为定位 DSL、数据隔离或执行器问题提供依据。"),
        ("评审重点", "项目已经进入工程闭环阶段，而非停留在静态原型或单次演示。")
    ], "复测记录体现了项目的工程成熟度：能失败、能分析、能修正、能再次验证。")

    add_image_slide(prs, 7, "功能证据 3：用例中心与测试资产管理", "AI 生成结果最终沉淀为项目级测试资产，支持后续执行和维护。", "cases", [
        ("资产结构", "当前平台采用 Project → Case → Run 结构，便于按项目管理用例和执行历史。"),
        ("业务示例", "login-brand-filter-cart 覆盖登录、品牌筛选、购物车、价格与数量断言。"),
        ("平台价值", "用例不是临时 prompt，而是可保存、可编辑、可复跑的自动化测试资产。")
    ], "测试资产沉淀是平台区别于“一次性脚本生成”的关键。", image_side="right")

    add_image_slide(prs, 8, "功能证据 4：结构化 DSL 编辑", "通过稳定动作集合表达测试流程，保证 AI 生成结果可审查、可修改、可执行。", "edit", [
        ("动作受控", "第一阶段动作集合保持小而稳定，例如 goto、click、input、wait_for、assert_text。"),
        ("变量复用", "支持 ${login_email}、${login_password} 等变量，占位数据可在执行时替换。"),
        ("执行可信", "所有正式测试执行都基于结构化 DSL，而不是直接把自然语言交给浏览器。")
    ], "DSL 是本项目控制 AI 不确定性的关键工艺层。")

    # 9 Locator innovation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "技术创新点：混合定位系统降低 UI 自动化脆弱性", "DOM 优先、视觉兜底、人工修正沉淀，而不是单一 CSS/XPath 或纯视觉 Agent。", "9/15")
    headers = ["层级", "策略", "作用", "创新价值"]
    rows = [
        ["Tier 0", "人工修正记录", "优先复用历史确认过的 selector", "把一次人工修复变成长期资产"],
        ["Tier 1", "DOM 语义定位", "基于 role、label、placeholder、text、aria 等召回候选", "成本低、速度快、可解释"],
        ["Tier 1.5/2", "上下文增强定位", "支持“商品附近按钮”等语义关系", "解决同名按钮和列表场景"],
        ["Tier 3", "AI visual 兜底", "截图交给 VLM 返回 bbox，再反查 DOM", "处理 DOM 信息不足的疑难场景"],
        ["Tier 4", "人工干预", "记录截图、URL、DOM、候选与失败原因", "为复盘和后续修正提供证据"],
    ]
    add_table(slide, headers, rows, 0.75, 1.55, 11.85, 4.45, [0.13, 0.20, 0.37, 0.30], 11)
    add_takeaway(slide, "相比纯脚本定位，它更能适应页面变化；相比纯视觉定位，它成本更低、解释性更强。")

    add_image_slide(prs, 10, "功能证据 5：步骤级证据报告", "每一步都有截图和执行上下文，失败后可以复盘、定位和归因。", "evidence", [
        ("证据内容", "步骤编号、动作类型、目标文本、耗时、截图、断言结果和页面状态。"),
        ("排障价值", "失败时可以看到发生在哪一步、页面实际长什么样、断言目标是否匹配。"),
        ("AI 分析基础", "结构化报告未来可输入 AI 失败分析模块，自动判断数据、定位、断言或环境问题。")
    ], "报告不是黑盒结果，而是可以被人和 AI 共同分析的数据资产。")

    add_image_slide(prs, 11, "功能证据 6：项目级报告中心", "从单次执行上升到项目维度质量观察。", "report", [
        ("当前能力", "项目级统计、执行列表、失败类型标签、历史执行记录和详情跳转。"),
        ("管理价值", "团队可以从执行历史中观察回归稳定性，而不只是查看单次脚本输出。"),
        ("扩展方向", "后续可加入趋势图、失败原因聚合、AI 失败分析和批量回归编排。")
    ], "报告中心把自动化测试从脚本工具提升为可运营的平台。", image_side="right")

    # 12 Application
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "应用场景与推广计划", "先服务高频 Web 回归，再沉淀行业模板和企业化能力。", "12/15")
    headers = ["阶段", "目标用户", "主推能力", "验证指标"]
    rows = [
        ["Demo 验证", "评委、同学、测试开发学习者", "AI 规划、用例执行、步骤证据报告", "能否独立看懂并复现完整链路"],
        ["团队试用", "中小研发团队、QA 团队", "项目级用例管理、核心路径回归", "脚本编写时间和回归时间下降"],
        ["行业模板", "电商、SaaS、后台系统团队", "登录、表单、购物车、筛选查询模板", "新项目自动化启动时间缩短"],
        ["企业化", "有内网系统和合规需求的企业", "私有化部署、权限、CI/CD、报告分析", "发布质量和排障效率提升"],
    ]
    add_table(slide, headers, rows, 0.75, 1.55, 11.85, 3.75, [0.18, 0.25, 0.32, 0.25], 11)
    add_takeaway(slide, "最适合切入的早期场景：Web 系统冒烟测试、发布前核心路径回归、测试资产沉淀。")

    # 13 Cost
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "成本数据分析：成本可控来自“DOM 优先 + AI 按需调用”", "评审估算口径，实际部署可按并发、模型和报告留存周期替换。", "13/15")
    headers = ["成本项", "估算口径", "控制方式"]
    rows = [
        ["研发成本", "1 名全栈/测试开发，3-6 个月完成 MVP 到可演示闭环", "优先打通 DSL → Runner → Evidence → Report 主链路"],
        ["服务器成本", "MVP 阶段 2C4G 云主机 + 数据库 + 对象存储，约 200-800 元/月", "低并发先单机部署，后续横向扩展执行器"],
        ["LLM 调用成本", "主要发生在规划、DSL 生成、疑难定位和失败分析", "常规执行不依赖 LLM，DOM 定位优先"],
        ["浏览器执行成本", "Playwright 执行消耗 CPU/内存，取决于并发数", "通过队列和并发上限控制资源"],
        ["存储成本", "截图、DOM 摘要、日志和 JSON 报告", "按项目和时间设置报告保留策略"],
    ]
    add_table(slide, headers, rows, 0.75, 1.55, 11.85, 4.35, [0.18, 0.47, 0.35], 11)
    add_takeaway(slide, "架构上避免“每一步都靠大模型”，因此比纯视觉 Agent 更适合长期自动化测试运行。")

    # 14 ROI
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header(slide, "投资回报与商业化分析", "价值来自减少重复回归、降低脚本维护和缩短失败排查时间。", "14/15")
    headers = ["维度", "示例/假设", "收益判断"]
    rows = [
        ["人工回归", "每月 4 次发布，每次核心回归 2 人天，按 800 元/人天估算，月度约 6400 元", "若覆盖 50%-70% 重复回归，每月可节省约 3200-4480 元"],
        ["脚本维护", "页面改版、同名按钮、动态内容导致定位失效", "混合定位和人工修正沉淀可减少重复维护"],
        ["失败排查", "传统报告只给 pass/fail，定位失败原因依赖人工复现", "步骤级证据让问题更快归因到数据、断言、定位或环境"],
        ["商业模式", "SaaS 订阅、私有化部署、测试服务套餐、CI/CD 集成", "适合从中小团队试用和测试服务交付切入"],
    ]
    add_table(slide, headers, rows, 0.75, 1.55, 11.85, 3.95, [0.18, 0.48, 0.34], 11)
    add_takeaway(slide, "当团队存在稳定 Web 产品、频繁发布和重复回归时，本项目具备清晰投入回报逻辑。")

    # 15 Summary
    add_image_slide(prs, 15, "总结：本项目的核心价值是“AI + 测试工程治理”", "AI 参与规划和分析，结构化 DSL 与后端 Runner 保证执行可信。", "sessions", [
        ("技术完整性", "AI Planning、DSL、Runner、Locator、Reporter 已形成端到端链路。"),
        ("创新合理性", "DOM 优先、视觉兜底、人工修正沉淀，比纯脚本和纯 Agent 更适合测试场景。"),
        ("商业可行性", "面向高频 Web 回归、自动化起步团队和测试服务交付，有明确成本节省空间。")
    ], "一句话：这不是让 AI 随机操作网页，而是把 AI 放进可审查、可执行、可复盘的自动化测试闭环。", image_side="right")

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
